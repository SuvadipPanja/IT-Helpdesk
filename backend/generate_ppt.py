import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()

# Path to screenshots
screenshot_dir = r"d:\Project\it-helpdesk\frontend\presentation_screenshots"

# Function to add a title slide
def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

# Function to add a content slide with bullets
def add_bullet_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for pt in bullet_points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(24)

# Function to add image slide
def add_image_slide(title_text, img_filename):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    img_path = os.path.join(screenshot_dir, img_filename)
    if os.path.exists(img_path):
        # Center the image roughly
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture(img_path, left, top, height=height)
    else:
        # fallback text
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        tf = txBox.text_frame
        tf.text = f"Image not found: {img_filename}"

# Create the slides
add_title_slide(
    "Nexus Support - IT Helpdesk Platform",
    "Comprehensive Overview of Architecture, Modules & AI Capabilities\nClient Presentation"
)

add_bullet_slide(
    "Robust System Architecture",
    [
        "Frontend: React JS with Vite for lightning-fast performance.",
        "Backend: Node.js (Express) ensuring high scalability and throughput.",
        "Database: Microsoft SQL Server (MSSQL) for enterprise-grade data integrity.",
        "Deployment ready: Fully Dockerized for seamless container orchestration."
    ]
)

add_bullet_slide(
    "Advanced Security & Licensing",
    [
        "Hardware-bound Offline Licensing Model.",
        "Generates unique fingerprints linked to physical machines.",
        "Requires zero internet connection for validation (Highly Secure).",
        "Role-Based Access Control (RBAC) separating Admins, Users, and Agents."
    ]
)

add_bullet_slide(
    "Intelligent AI & Bot Engine",
    [
        "Local Ollama Module: AI capabilities run 100% locally.",
        "Self-Learning Context Awareness for resolving tickets autonomously.",
        "Spelling Correction & Intent detection services.",
        "Bot Training Service & API Integration to customize behavior."
    ]
)

add_bullet_slide(
    "Master Control & Modules",
    [
        "Analytics Module: Real-time insights into SLAs and agent performance.",
        "Settings Module: Centralized master control for departments, roles, and UI.",
        "Ticket Module: Auto-assignment, comprehensive tracking, and attachments.",
        "Job Monitor: UI for admins to start, stop, and track background jobs."
    ]
)

# Adding screenshot slides
add_image_slide("Administrator Master Control - Dashboard", "Admin_dashboard.png")
add_image_slide("Comprehensive Analytics", "Admin_analytics.png")
add_image_slide("Settings & Master Configurations", "Admin_settings.png")
add_image_slide("Hardware-Bound License Management", "Admin_license.png")
add_image_slide("Administrator Ticket Management", "Admin_tickets.png")
add_image_slide("User Experience - Dashboard", "User_dashboard.png")
add_image_slide("User Experience - Support Tickets", "User_tickets.png")

# Save
output_path = r"d:\Project\it-helpdesk\frontend\IT_Helpdesk_Presentation.pptx"
prs.save(output_path)
print(f"Presentation generated at: {output_path}")
